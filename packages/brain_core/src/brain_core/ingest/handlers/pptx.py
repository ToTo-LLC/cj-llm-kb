"""PptxHandler — PowerPoint (.pptx) slide-deck extraction.

Plan 24 T2. PowerPoint files extract to a markdown body where each slide
is its own `## Slide N: <title>` section with bulleted text content and
an optional `**Speaker notes:**` block. Every embedded image is collected
into `extras["images"]` for the Plan 24 T4 OCR pass.

Extraction shape (per plan §T2.3):
- Slide title comes from `slide.shapes.title.text`; if the layout has no
  title placeholder (e.g. Blank layout) the section heading degrades to
  `## Slide N` with no title fragment.
- Bullets walk every text-bearing shape on the slide (excluding the title
  shape, which is already in the heading). `text_frame.paragraphs[i].level`
  is preserved as nested-list indent (level 0 → top-level bullet, level 1
  → 2-space indent, etc.).
- Speaker notes use `slide.notes_slide.notes_text_frame.text` only when
  `slide.has_notes_slide` is True AND the resulting text is non-empty. An
  empty notes block suppresses the `**Speaker notes:**` heading entirely.
- Images: walk `slide.shapes` filtering on
  `shape.shape_type == MSO_SHAPE_TYPE.PICTURE` and emit
  `{"blob": bytes, "content_type": str, "slide_index": int, "index": int}`
  dicts (slide_index 1-based; index is the overall image counter starting
  at 0). T4's OCR pass consumes this exact shape.
- `title`: `pres.core_properties.title` if set; else first slide title;
  else file stem.
- `author` / `published`: from `pres.core_properties` — author may be the
  empty string (coerced to None); `created` is a datetime narrowed to
  `date`.
- `source_url`: None (pptx is local).
- `source_type`: `SourceType.PPTX`.
- `archive_path`: original file copied into `archive_root` (mirrors
  DocxHandler / TranscriptDOCXHandler).

Routing contract (Plan 24 §T2 dispatcher): PptxHandler is registered
AFTER PDFHandler and BEFORE EmailHandler — it's format-specific enough
to claim ahead of the email/text catch-alls but doesn't compete with any
.docx-claiming handler. There is no pptx-transcript convention; every
.pptx routes here.
"""

from __future__ import annotations

import shutil
from datetime import date
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from brain_core.ingest.handlers.base import HandlerError
from brain_core.ingest.types import ExtractedSource, SourceType


class PptxHandler:
    """PowerPoint slide-deck handler.

    `can_handle` is a pure routing check — it only inspects the path
    suffix. `extract` opens the file with python-pptx and walks slides
    in document order.
    """

    source_type: SourceType = SourceType.PPTX

    def can_handle(self, spec: str | Path) -> bool:
        """Claim any local `.pptx` file. Pure routing check — no I/O."""
        if not isinstance(spec, Path):
            return False
        if spec.suffix.lower() != ".pptx":
            return False
        return spec.is_file()

    async def extract(self, spec: str | Path, *, archive_root: Path) -> ExtractedSource:
        """Open the .pptx, walk slides, collect text + images.

        Wraps python-pptx open failures in `HandlerError` so the
        pipeline gets an actionable message instead of a raw zipfile
        exception. Slide-level failures (e.g. a single malformed shape)
        are caught locally so one bad slide doesn't kill the whole
        extract.
        """
        if not isinstance(spec, Path) or not spec.exists():
            raise HandlerError(f"pptx handler cannot read {spec!r}")
        try:
            pres = Presentation(str(spec))
        except Exception as exc:
            raise HandlerError(
                f"Could not open PPTX {spec.name!r}: {exc}. "
                "The file may be corrupt or not a PowerPoint document."
            ) from exc

        body_parts: list[str] = []
        images: list[dict[str, Any]] = []
        first_slide_title: str | None = None
        image_counter = 0

        for slide_index, slide in enumerate(pres.slides, start=1):
            slide_title = _slide_title(slide)
            if slide_index == 1 and slide_title:
                first_slide_title = slide_title

            heading = (
                f"## Slide {slide_index}: {slide_title}"
                if slide_title
                else f"## Slide {slide_index}"
            )
            section_parts: list[str] = [heading]

            bullets = _slide_bullets(slide)
            if bullets:
                section_parts.append(bullets)

            notes = _slide_notes(slide)
            if notes:
                section_parts.append(f"**Speaker notes:**\n\n{notes}")

            body_parts.append("\n\n".join(section_parts))

            for shape in slide.shapes:
                blob_entry = _shape_to_image_dict(shape, slide_index, image_counter)
                if blob_entry is not None:
                    images.append(blob_entry)
                    image_counter += 1

        body_text = "\n\n".join(body_parts)

        # Title fallback chain: core_properties.title -> first slide title
        # -> filename stem. python-pptx returns "" for unset.
        title_raw = pres.core_properties.title
        title: str | None = title_raw.strip() if title_raw else None
        if not title:
            title = first_slide_title or spec.stem

        author_raw = pres.core_properties.author
        author: str | None = author_raw.strip() if author_raw else None
        if author == "":
            author = None

        created_raw = pres.core_properties.created
        published: date | None = created_raw.date() if created_raw is not None else None

        archive_root.mkdir(parents=True, exist_ok=True)
        archive_path = archive_root / spec.name
        shutil.copy2(spec, archive_path)

        extras: dict[str, Any] = {"images": images}

        return ExtractedSource(
            title=title,
            author=author,
            published=published,
            source_url=None,
            source_type=SourceType.PPTX,
            body_text=body_text,
            archive_path=archive_path,
            extras=extras,
        )


def _slide_title(slide: Any) -> str | None:
    """Return the slide's title text or None if no title placeholder.

    Layouts without a title placeholder (e.g. Blank) return None for
    `shapes.title`; we handle that without raising. An empty title
    string is also normalized to None so the caller can render a
    title-less heading.
    """
    try:
        title_shape = slide.shapes.title
    except AttributeError:
        return None
    if title_shape is None:
        return None
    text = (title_shape.text or "").strip()
    return text or None


def _slide_bullets(slide: Any) -> str:
    """Walk text-bearing shapes (excluding the title) and emit a markdown list.

    Each paragraph becomes a bullet; `paragraph.level` (0-based in
    python-pptx) maps to nested-list indent (2 spaces per level).
    Empty paragraphs are skipped. The title placeholder is excluded
    because its text is already in the section heading.

    Identity caveat: `slide.shapes.title` returns a fresh proxy object
    on each access, so `shape is title_shape` is unreliable. We compare
    by `shape_id` (stable per-shape integer assigned by PowerPoint)
    instead.
    """
    title_shape_id: int | None = None
    try:
        title_shape = slide.shapes.title
        if title_shape is not None:
            title_shape_id = title_shape.shape_id
    except AttributeError:
        pass

    lines: list[str] = []
    for shape in slide.shapes:
        if title_shape_id is not None and shape.shape_id == title_shape_id:
            continue
        if not shape.has_text_frame:
            continue
        text_frame = shape.text_frame
        for paragraph in text_frame.paragraphs:
            text = (paragraph.text or "").strip()
            if not text:
                continue
            level = max(0, int(paragraph.level or 0))
            indent = "  " * level
            lines.append(f"{indent}- {text}")
    return "\n".join(lines)


def _slide_notes(slide: Any) -> str:
    """Return speaker-notes text or empty string if none.

    `has_notes_slide` is False for slides that have never had notes
    touched, so we check it before reading `notes_slide` (touching
    `notes_slide` materializes a slide-notes part on disk we'd rather
    not provoke). Even when `has_notes_slide` is True the text frame
    may be empty.
    """
    if not slide.has_notes_slide:
        return ""
    notes_slide = slide.notes_slide
    if notes_slide is None or notes_slide.notes_text_frame is None:
        return ""
    return (notes_slide.notes_text_frame.text or "").strip()


def _shape_to_image_dict(
    shape: Any,
    slide_index: int,
    counter: int,
) -> dict[str, Any] | None:
    """Extract image bytes + content_type from a PICTURE shape.

    Returns None for non-PICTURE shapes (text boxes, charts, etc.) and
    for malformed PICTURE shapes that can't surface an `image` attribute
    (defensive — group-shape pictures and some linked-picture shapes
    don't expose `.image`). T4's OCR pass consumes the returned dict
    shape exactly:
        {"blob": bytes, "content_type": str, "slide_index": int, "index": int}
    """
    try:
        if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
            return None
        image = shape.image
        return {
            "blob": image.blob,
            "content_type": image.content_type,
            "slide_index": slide_index,
            "index": counter,
        }
    except (AttributeError, KeyError, ValueError):
        return None

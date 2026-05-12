"""Unit tests for PptxHandler (Plan 24 T2).

Fixtures are constructed at runtime via python-pptx so the repo stays
free of binary .pptx files. Mirrors the test_docx.py pattern.
"""

from __future__ import annotations

import io
import struct
import zlib
from pathlib import Path

import pytest
from brain_core.ingest.dispatcher import _default_handlers, dispatch
from brain_core.ingest.handlers.base import HandlerError
from brain_core.ingest.handlers.email import EmailHandler
from brain_core.ingest.handlers.pdf import PDFHandler
from brain_core.ingest.handlers.pptx import PptxHandler
from brain_core.ingest.types import SourceType
from pptx import Presentation
from pptx.util import Inches

# ---------------------------------------------------------------------------
# Test helpers
# ---------------------------------------------------------------------------


def _make_tiny_png() -> bytes:
    """Build a valid 1x1 white RGB PNG via stdlib zlib + struct.

    Avoids PIL / binary fixture files so the test is hermetic.
    Same shape as the helper in test_docx.py.
    """
    sig = b"\x89PNG\r\n\x1a\n"
    ihdr_data = struct.pack(">IIBBBBB", 1, 1, 8, 2, 0, 0, 0)
    ihdr = (
        struct.pack(">I", 13)
        + b"IHDR"
        + ihdr_data
        + struct.pack(">I", zlib.crc32(b"IHDR" + ihdr_data))
    )
    raw = b"\x00\xff\xff\xff"
    compressed = zlib.compress(raw)
    idat = (
        struct.pack(">I", len(compressed))
        + b"IDAT"
        + compressed
        + struct.pack(">I", zlib.crc32(b"IDAT" + compressed))
    )
    iend = (
        struct.pack(">I", 0)
        + b"IEND"
        + struct.pack(">I", zlib.crc32(b"IEND"))
    )
    return sig + ihdr + idat + iend


def _build_three_slide_pptx(tmp_path: Path, name: str = "deck.pptx") -> Path:
    """Three-slide deck: each slide has a title + bulleted content."""
    pres = Presentation()
    layout = pres.slide_layouts[1]  # Title and Content

    s1 = pres.slides.add_slide(layout)
    s1.shapes.title.text = "Strategy 2026"
    s1.placeholders[1].text = "Bullet one\nBullet two"

    s2 = pres.slides.add_slide(layout)
    s2.shapes.title.text = "Q1 Themes"
    s2.placeholders[1].text = "Theme A\nTheme B"

    s3 = pres.slides.add_slide(layout)
    s3.shapes.title.text = "Roadmap"
    s3.placeholders[1].text = "Milestone 1"

    path = tmp_path / name
    pres.save(str(path))
    return path


def _build_pptx_with_bullets(tmp_path: Path, name: str = "bullets.pptx") -> Path:
    pres = Presentation()
    slide = pres.slides.add_slide(pres.slide_layouts[1])
    slide.shapes.title.text = "Bullet List"
    slide.placeholders[1].text = "Alpha\nBravo\nCharlie"
    path = tmp_path / name
    pres.save(str(path))
    return path


def _build_pptx_with_notes(tmp_path: Path, name: str = "with-notes.pptx") -> Path:
    pres = Presentation()
    slide = pres.slides.add_slide(pres.slide_layouts[1])
    slide.shapes.title.text = "Slide With Notes"
    slide.placeholders[1].text = "Body bullet"
    slide.notes_slide.notes_text_frame.text = "Remember to mention the pivot."
    path = tmp_path / name
    pres.save(str(path))
    return path


def _build_pptx_without_notes(tmp_path: Path, name: str = "no-notes.pptx") -> Path:
    pres = Presentation()
    slide = pres.slides.add_slide(pres.slide_layouts[1])
    slide.shapes.title.text = "Notes-Free Slide"
    slide.placeholders[1].text = "Body bullet"
    # Do NOT touch slide.notes_slide — keeps has_notes_slide False.
    path = tmp_path / name
    pres.save(str(path))
    return path


def _build_pptx_with_image(tmp_path: Path, name: str = "with-image.pptx") -> Path:
    pres = Presentation()
    slide = pres.slides.add_slide(pres.slide_layouts[1])
    slide.shapes.title.text = "Image Slide"
    slide.placeholders[1].text = "Caption text"
    slide.shapes.add_picture(io.BytesIO(_make_tiny_png()), Inches(2), Inches(2))
    path = tmp_path / name
    pres.save(str(path))
    return path


def _build_pptx_image_only(tmp_path: Path, name: str = "image-only.pptx") -> Path:
    """Slide using the Blank layout with only an image — no title, no bullets."""
    pres = Presentation()
    slide = pres.slides.add_slide(pres.slide_layouts[6])  # Blank
    slide.shapes.add_picture(io.BytesIO(_make_tiny_png()), Inches(2), Inches(2))
    path = tmp_path / name
    pres.save(str(path))
    return path


def _build_pptx_missing_title(tmp_path: Path, name: str = "no-title.pptx") -> Path:
    """First slide is a Blank layout — no title placeholder available."""
    pres = Presentation()
    pres.slides.add_slide(pres.slide_layouts[6])  # Blank
    path = tmp_path / name
    pres.save(str(path))
    return path


def _build_pptx_with_core_title(
    tmp_path: Path,
    name: str = "titled.pptx",
    core_title: str = "Workshop Notes 2026",
) -> Path:
    pres = Presentation()
    pres.core_properties.title = core_title
    slide = pres.slides.add_slide(pres.slide_layouts[1])
    slide.shapes.title.text = "Different slide title"
    path = tmp_path / name
    pres.save(str(path))
    return path


# ---------------------------------------------------------------------------
# Routing
# ---------------------------------------------------------------------------


def test_can_handle_claims_pptx_suffix(tmp_path: Path) -> None:
    """Suffix-only routing: .pptx claimed; .docx/.pdf/.txt rejected."""
    h = PptxHandler()
    pptx = _build_three_slide_pptx(tmp_path)
    assert h.can_handle(pptx) is True
    # Case-insensitive suffix match
    upper = tmp_path / "shouty.PPTX"
    upper.write_bytes(pptx.read_bytes())
    assert h.can_handle(upper) is True
    # Other suffixes rejected
    docx_path = tmp_path / "foo.docx"
    docx_path.write_bytes(b"not even a docx")
    pdf_path = tmp_path / "foo.pdf"
    pdf_path.write_bytes(b"%PDF-1.4")
    txt_path = tmp_path / "foo.txt"
    txt_path.write_text("plain text", encoding="utf-8")
    assert h.can_handle(docx_path) is False
    assert h.can_handle(pdf_path) is False
    assert h.can_handle(txt_path) is False
    # Non-Path inputs and non-existent paths are not claimed
    assert h.can_handle("https://example.com/deck.pptx") is False
    assert h.can_handle(tmp_path / "missing.pptx") is False


def test_dispatcher_registers_pptx_between_pdf_and_email() -> None:
    """Plan 24 §T2 pins PptxHandler between PDFHandler and EmailHandler."""
    chain = _default_handlers()
    pdf_idx = next(i for i, h in enumerate(chain) if isinstance(h, PDFHandler))
    pptx_idx = next(i for i, h in enumerate(chain) if isinstance(h, PptxHandler))
    email_idx = next(i for i, h in enumerate(chain) if isinstance(h, EmailHandler))
    assert pdf_idx < pptx_idx < email_idx


@pytest.mark.asyncio
async def test_dispatcher_routes_pptx_to_pptx_handler(tmp_path: Path) -> None:
    path = _build_three_slide_pptx(tmp_path, name="strategy.pptx")
    handler = await dispatch(path)
    assert isinstance(handler, PptxHandler)


# ---------------------------------------------------------------------------
# Extraction
# ---------------------------------------------------------------------------


@pytest.mark.asyncio
async def test_extract_slides_to_markdown_sections(tmp_path: Path) -> None:
    """3-slide fixture emits 3 ## Slide N: headings in document order."""
    path = _build_three_slide_pptx(tmp_path)
    es = await PptxHandler().extract(path, archive_root=tmp_path / "archive")
    assert es.source_type is SourceType.PPTX
    assert "## Slide 1: Strategy 2026" in es.body_text
    assert "## Slide 2: Q1 Themes" in es.body_text
    assert "## Slide 3: Roadmap" in es.body_text
    # Order preserved
    idx1 = es.body_text.index("## Slide 1")
    idx2 = es.body_text.index("## Slide 2")
    idx3 = es.body_text.index("## Slide 3")
    assert idx1 < idx2 < idx3


@pytest.mark.asyncio
async def test_extract_slide_titles(tmp_path: Path) -> None:
    """Slide titles appear verbatim in their section headings."""
    path = _build_three_slide_pptx(tmp_path)
    es = await PptxHandler().extract(path, archive_root=tmp_path / "archive")
    for title in ("Strategy 2026", "Q1 Themes", "Roadmap"):
        assert title in es.body_text


@pytest.mark.asyncio
async def test_extract_bullet_points_to_markdown(tmp_path: Path) -> None:
    """Bulleted content frames render as markdown list items."""
    path = _build_pptx_with_bullets(tmp_path)
    es = await PptxHandler().extract(path, archive_root=tmp_path / "archive")
    assert "- Alpha" in es.body_text
    assert "- Bravo" in es.body_text
    assert "- Charlie" in es.body_text
    # The title line itself is NOT bulleted — it's the section heading.
    assert "- Bullet List" not in es.body_text


@pytest.mark.asyncio
async def test_extract_speaker_notes(tmp_path: Path) -> None:
    """Slides with notes emit a **Speaker notes:** block carrying the text."""
    path = _build_pptx_with_notes(tmp_path)
    es = await PptxHandler().extract(path, archive_root=tmp_path / "archive")
    assert "**Speaker notes:**" in es.body_text
    assert "Remember to mention the pivot." in es.body_text


@pytest.mark.asyncio
async def test_extract_skips_empty_speaker_notes(tmp_path: Path) -> None:
    """Slides without notes do NOT emit the **Speaker notes:** heading."""
    path = _build_pptx_without_notes(tmp_path)
    es = await PptxHandler().extract(path, archive_root=tmp_path / "archive")
    assert "**Speaker notes:**" not in es.body_text
    # But the slide title + body still render correctly.
    assert "Notes-Free Slide" in es.body_text
    assert "- Body bullet" in es.body_text


@pytest.mark.asyncio
async def test_extract_collects_slide_images(tmp_path: Path) -> None:
    """Embedded slide pictures land in extras["images"] with the T4 shape."""
    path = _build_pptx_with_image(tmp_path)
    es = await PptxHandler().extract(path, archive_root=tmp_path / "archive")
    images = es.extras.get("images", [])
    assert len(images) == 1
    img = images[0]
    # T4's OCR pass expects this exact shape:
    assert isinstance(img["blob"], bytes)
    assert len(img["blob"]) > 0
    assert img["content_type"] == "image/png"
    assert img["slide_index"] == 1
    assert img["index"] == 0


@pytest.mark.asyncio
async def test_extract_handles_image_only_slide(tmp_path: Path) -> None:
    """Slide with no text shapes still renders heading + image collected."""
    path = _build_pptx_image_only(tmp_path)
    es = await PptxHandler().extract(path, archive_root=tmp_path / "archive")
    # Heading present (no title since Blank layout has none)
    assert "## Slide 1" in es.body_text
    # No bullets since there are no text shapes
    assert "- " not in es.body_text
    # No speaker-notes section since notes_slide wasn't touched
    assert "**Speaker notes:**" not in es.body_text
    images = es.extras.get("images", [])
    assert len(images) == 1
    assert images[0]["slide_index"] == 1


@pytest.mark.asyncio
async def test_extract_handles_missing_title(tmp_path: Path) -> None:
    """Blank-layout first slide has no title placeholder — no crash, no title fragment."""
    path = _build_pptx_missing_title(tmp_path)
    es = await PptxHandler().extract(path, archive_root=tmp_path / "archive")
    # Section heading uses the title-less form
    assert "## Slide 1" in es.body_text
    assert "## Slide 1:" not in es.body_text  # No colon-prefixed title fragment
    # ExtractedSource.title falls back to filename stem when no title anywhere
    assert es.title == "no-title"


@pytest.mark.asyncio
async def test_extract_uses_core_properties_title(tmp_path: Path) -> None:
    """pres.core_properties.title wins over slide title for ExtractedSource.title."""
    path = _build_pptx_with_core_title(tmp_path, core_title="Workshop Notes 2026")
    es = await PptxHandler().extract(path, archive_root=tmp_path / "archive")
    assert es.title == "Workshop Notes 2026"


# ---------------------------------------------------------------------------
# Error paths
# ---------------------------------------------------------------------------


@pytest.mark.asyncio
async def test_extract_raises_handler_error_on_corrupt_pptx(tmp_path: Path) -> None:
    """Garbage with .pptx extension raises HandlerError, not a raw zipfile crash."""
    fake = tmp_path / "fake.pptx"
    fake.write_bytes(b"not a real pptx file")
    with pytest.raises(HandlerError):
        await PptxHandler().extract(fake, archive_root=tmp_path / "archive")


@pytest.mark.asyncio
async def test_extract_raises_when_file_missing(tmp_path: Path) -> None:
    missing = tmp_path / "missing.pptx"
    with pytest.raises(HandlerError):
        await PptxHandler().extract(missing, archive_root=tmp_path / "archive")


@pytest.mark.asyncio
async def test_archive_path_copies_file(tmp_path: Path) -> None:
    """ExtractedSource.archive_path points at a fresh copy of the original."""
    path = _build_three_slide_pptx(tmp_path)
    archive = tmp_path / "archive"
    es = await PptxHandler().extract(path, archive_root=archive)
    assert es.archive_path.exists()
    assert es.archive_path.parent == archive
    # Byte-identical copy
    assert es.archive_path.read_bytes() == path.read_bytes()

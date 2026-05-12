"""Plan 24 Task 4 — integration tests for the pipeline OCR pass on PptxHandler output.

End-to-end coverage: build .pptx fixtures with images on specific slides, run
them through :class:`IngestPipeline` with :class:`FakeLLMProvider` priming the
summarize / integrate / vision queues, and assert the slide-prefixed inline
blocks land in the OCR'd ``body_text`` returned on ``IngestResult.extracted``.

These tests pin the PPTX-specific inline format
(``[Image (slide N): <text>]``) and per-image error isolation (one bad image
on slide K doesn't kill the OCR pass for the others).
"""

from __future__ import annotations

import io
import struct
import zlib
from datetime import UTC, datetime
from pathlib import Path

import pytest
from brain_core.budget import PerDomainBudgetGuard
from brain_core.config.schema import Config
from brain_core.cost.ledger import CostLedger
from brain_core.ingest.pipeline import IngestPipeline
from brain_core.ingest.types import IngestStatus
from brain_core.llm.fake import FakeLLMProvider
from brain_core.prompts.schemas import SummarizeOutput
from brain_core.vault.types import IndexEntryPatch, PatchSet
from brain_core.vault.writer import VaultWriter
from pptx import Presentation
from pptx.util import Inches
from structlog.testing import capture_logs

# ---------------------------------------------------------------------------
# Hermetic fixture builders
# ---------------------------------------------------------------------------


def _make_tiny_png() -> bytes:
    """1x1 white RGB PNG — same helper as ``test_pptx.py``."""
    sig = b"\x89PNG\r\n\x1a\n"
    ihdr_data = struct.pack(">IIBBBBB", 1, 1, 8, 2, 0, 0, 0)
    ihdr = (
        struct.pack(">I", 13)
        + b"IHDR"
        + ihdr_data
        + struct.pack(">I", zlib.crc32(b"IHDR" + ihdr_data))
    )
    raw = b"\x00\xff\xff\xff"
    compressed = zlib.compress(raw)
    idat = (
        struct.pack(">I", len(compressed))
        + b"IDAT"
        + compressed
        + struct.pack(">I", zlib.crc32(b"IDAT" + compressed))
    )
    iend = (
        struct.pack(">I", 0)
        + b"IEND"
        + struct.pack(">I", zlib.crc32(b"IEND"))
    )
    return sig + ihdr + idat + iend


def _build_three_slide_pptx_with_image_on_slide_2(tmp_path: Path) -> Path:
    """3 slides; image on slide 2 only."""
    pres = Presentation()
    layout = pres.slide_layouts[1]  # Title and Content

    s1 = pres.slides.add_slide(layout)
    s1.shapes.title.text = "Cover"
    s1.placeholders[1].text = "intro bullet"

    s2 = pres.slides.add_slide(layout)
    s2.shapes.title.text = "Chart slide"
    s2.placeholders[1].text = "context"
    s2.shapes.add_picture(io.BytesIO(_make_tiny_png()), Inches(2), Inches(2))

    s3 = pres.slides.add_slide(layout)
    s3.shapes.title.text = "Roadmap"
    s3.placeholders[1].text = "milestones"

    path = tmp_path / "deck.pptx"
    pres.save(str(path))
    return path


def _build_five_slide_pptx_with_images_on_2_and_5(tmp_path: Path) -> Path:
    """5 slides; images on slides 2 and 5."""
    pres = Presentation()
    layout = pres.slide_layouts[1]

    titles = ["Cover", "Image two", "Body", "Body", "Image five"]
    for idx, title in enumerate(titles, start=1):
        s = pres.slides.add_slide(layout)
        s.shapes.title.text = title
        s.placeholders[1].text = f"content for slide {idx}"
        if idx in (2, 5):
            s.shapes.add_picture(io.BytesIO(_make_tiny_png()), Inches(2), Inches(2))

    path = tmp_path / "multi-image-deck.pptx"
    pres.save(str(path))
    return path


def _build_plain_pptx(tmp_path: Path) -> Path:
    """3 slides; zero images. Pins the no-OCR-call path."""
    pres = Presentation()
    layout = pres.slide_layouts[1]
    for title in ("Cover", "Middle", "End"):
        s = pres.slides.add_slide(layout)
        s.shapes.title.text = title
        s.placeholders[1].text = "body"
    path = tmp_path / "plain.pptx"
    pres.save(str(path))
    return path


def _build_two_image_pptx(tmp_path: Path) -> Path:
    """2 slides each carrying one image — used to pin partial-failure isolation."""
    pres = Presentation()
    layout = pres.slide_layouts[1]
    for idx, title in enumerate(("First image slide", "Second image slide"), start=1):
        s = pres.slides.add_slide(layout)
        s.shapes.title.text = title
        s.placeholders[1].text = f"slide {idx} body"
        s.shapes.add_picture(io.BytesIO(_make_tiny_png()), Inches(2), Inches(2))
    path = tmp_path / "two-image-deck.pptx"
    pres.save(str(path))
    return path


# ---------------------------------------------------------------------------
# Pipeline plumbing
# ---------------------------------------------------------------------------


def _queue_pipeline_responses(fake: FakeLLMProvider, *, title: str) -> None:
    """Queue summarize + integrate (classify skipped via domain_override)."""
    fake.queue(
        SummarizeOutput(
            title=title,
            summary="summary text",
            key_points=["point"],
            entities=[],
            concepts=[],
            open_questions=[],
        ).model_dump_json()
    )
    fake.queue(
        PatchSet(
            index_entries=[
                IndexEntryPatch(
                    section="Sources",
                    line=f"- [[{title}]] — summary text",
                    domain="research",
                )
            ],
            log_entry=f"## [2026-05-12 10:00] ingest | source | [[{title}]]",
            reason="initial",
        ).model_dump_json()
    )


def _make_pipeline(
    fake: FakeLLMProvider,
    *,
    vault: Path,
    cost_ledger: CostLedger,
    config: Config,
) -> IngestPipeline:
    guard = PerDomainBudgetGuard(cost_ledger)
    return IngestPipeline(
        vault_root=vault,
        writer=VaultWriter(vault_root=vault),
        llm=fake,
        summarize_model="claude-sonnet-4-6",
        integrate_model="claude-sonnet-4-6",
        classify_model="claude-haiku-4-5-20251001",
        guard=guard,
        config=config,
        cost_ledger=cost_ledger,
    )


# ---------------------------------------------------------------------------
# Tests
# ---------------------------------------------------------------------------


@pytest.mark.asyncio
async def test_pptx_with_slide_image_gets_slide_prefixed_block(
    ephemeral_vault: Path, tmp_path: Path
) -> None:
    """Image on slide 2 → ``[Image (slide 2): <text>]`` block appended to body."""
    pptx_path = _build_three_slide_pptx_with_image_on_slide_2(tmp_path)
    ledger = CostLedger(db_path=tmp_path / "costs.sqlite")
    config = Config(vault_path=ephemeral_vault)

    fake = FakeLLMProvider()
    fake.queue_vision("chart Q4 revenue")
    _queue_pipeline_responses(fake, title="chart-deck")

    p = _make_pipeline(fake, vault=ephemeral_vault, cost_ledger=ledger, config=config)

    r = await p.ingest(
        pptx_path,
        allowed_domains=("research",),
        domain_override="research",
    )

    assert r.status is IngestStatus.OK
    assert r.extracted is not None
    body = r.extracted.body_text
    assert "[Image (slide 2): chart Q4 revenue]" in body
    # Original slide structure preserved.
    assert "## Slide 1: Cover" in body
    assert "## Slide 2: Chart slide" in body
    assert "## Slide 3: Roadmap" in body
    # One vision call.
    assert len(fake.vision_calls) == 1
    # The summarize prompt saw the OCR'd body (regression pin).
    summarize_req = fake.requests[0]
    assert "[Image (slide 2): chart Q4 revenue]" in summarize_req.messages[0].content
    # Ledger booked the op="ocr" row to the research domain.
    today = datetime.now(tz=UTC).date()
    by_domain = ledger.total_by_domain(today)
    assert by_domain.get("research", 0.0) > 0.0


@pytest.mark.asyncio
async def test_pptx_with_multiple_slide_images_each_block_correct_number(
    ephemeral_vault: Path, tmp_path: Path
) -> None:
    """Images on slides 2 + 5 → both blocks present, each with the right slide number."""
    pptx_path = _build_five_slide_pptx_with_images_on_2_and_5(tmp_path)
    ledger = CostLedger(db_path=tmp_path / "costs.sqlite")
    config = Config(vault_path=ephemeral_vault)

    fake = FakeLLMProvider()
    # PptxHandler emits images in slide order, so queue OCR responses in
    # the same order: slide-2 text first, slide-5 text second.
    fake.queue_vision("slide two image text")
    fake.queue_vision("slide five image text")
    _queue_pipeline_responses(fake, title="multi-image")

    p = _make_pipeline(fake, vault=ephemeral_vault, cost_ledger=ledger, config=config)

    r = await p.ingest(
        pptx_path,
        allowed_domains=("research",),
        domain_override="research",
    )

    assert r.status is IngestStatus.OK
    assert r.extracted is not None
    body = r.extracted.body_text
    assert "[Image (slide 2): slide two image text]" in body
    assert "[Image (slide 5): slide five image text]" in body
    # Slide-2 block precedes slide-5 block (handler extraction order is
    # preserved through the OCR pass).
    assert body.index("(slide 2)") < body.index("(slide 5)")
    assert len(fake.vision_calls) == 2


@pytest.mark.asyncio
async def test_pptx_with_no_images_no_ocr_call(
    ephemeral_vault: Path, tmp_path: Path
) -> None:
    """Pure-text deck → vision_extract NEVER called.

    Vision queue stays empty (deliberate negative evidence).
    """
    pptx_path = _build_plain_pptx(tmp_path)
    ledger = CostLedger(db_path=tmp_path / "costs.sqlite")
    config = Config(vault_path=ephemeral_vault)

    fake = FakeLLMProvider()
    # NO vision queue. Empty-queue raise would surface if OCR fired.
    _queue_pipeline_responses(fake, title="plain-deck")

    p = _make_pipeline(fake, vault=ephemeral_vault, cost_ledger=ledger, config=config)

    r = await p.ingest(
        pptx_path,
        allowed_domains=("research",),
        domain_override="research",
    )

    assert r.status is IngestStatus.OK
    assert fake.vision_calls == []
    assert r.extracted is not None
    assert "[Image" not in r.extracted.body_text


@pytest.mark.asyncio
async def test_pptx_with_vision_error_on_second_image_first_block_survives(
    ephemeral_vault: Path, tmp_path: Path
) -> None:
    """2 images; first vision call succeeds, second raises (queue empty).

    Assert:
    * Ingest still succeeds (status=OK).
    * Body has exactly ONE inline block — for the first image.
    * A warning is logged for the second image.
    """
    pptx_path = _build_two_image_pptx(tmp_path)
    ledger = CostLedger(db_path=tmp_path / "costs.sqlite")
    config = Config(vault_path=ephemeral_vault)

    fake = FakeLLMProvider()
    # Queue exactly ONE vision response. The second OCR call drains the
    # queue and raises RuntimeError — the OCR pass catches it.
    fake.queue_vision("good ocr text")
    _queue_pipeline_responses(fake, title="partial-fail-deck")

    p = _make_pipeline(fake, vault=ephemeral_vault, cost_ledger=ledger, config=config)

    with capture_logs() as logs:
        r = await p.ingest(
            pptx_path,
            allowed_domains=("research",),
            domain_override="research",
        )

    assert r.status is IngestStatus.OK
    assert r.extracted is not None
    body = r.extracted.body_text
    assert "[Image (slide 1): good ocr text]" in body
    # Slide-2 block absent (the second image errored out).
    assert "[Image (slide 2):" not in body
    # Exactly one warning.
    warnings = [e for e in logs if e.get("event") == "ingest.ocr.image_skipped"]
    assert len(warnings) == 1
    assert warnings[0].get("slide_index") == 2
    # Both vision calls were attempted (the second one raised + was
    # caught — vision_calls is appended BEFORE the empty-queue raise).
    assert len(fake.vision_calls) == 2
